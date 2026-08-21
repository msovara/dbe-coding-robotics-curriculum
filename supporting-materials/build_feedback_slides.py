"""Build a short facilitator slide deck for the workshop quality feedback session."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Workshop_Quality_Feedback_Session.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x0D, 0x94, 0x88)
DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_run(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title_bar(slide, title):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title
    set_run(run, size=28, bold=True, color=WHITE)


def add_bullets(text_frame, items, size=18):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=DARK)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — title / how to run
    s1 = blank_slide(prs)
    add_title_bar(s1, "Workshop Quality Feedback Session")
    box = s1.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        ("Purpose", True, 22, TEAL),
        ("Get honest feedback on the quality of this workshop so the next cohort is better.", False, 18, DARK),
        ("", False, 12, DARK),
        ("How we will run this (15–20 minutes)", True, 22, TEAL),
        ("1. Frame (1 min) — honest feedback; names optional", False, 18, DARK),
        ("2. Silent write (5 min) — sticky notes or evaluation form", False, 18, DARK),
        ("3. Share-out (8–10 min) — each table: 1 keep + 1 improve", False, 18, DARK),
        ("4. Close (2 min) — thank you; note what we will act on", False, 18, DARK),
        ("", False, 12, DARK),
        ("If we only have 10 minutes: use questions 1, 2, 3, 8 and 10.", False, 16, MUTED),
    ]
    first = True
    for text, bold, size, color in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        if not text:
            continue
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=bold, color=color)

    # Slide 2 — questions
    s2 = blank_slide(prs)
    add_title_bar(s2, "Questions — Workshop Quality")
    left = s2.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(6.2), Inches(5.8))
    add_bullets(
        left.text_frame,
        [
            "1. Overall, how would you rate this workshop? (1–5)",
            "2. What was the strongest part of the week?",
            "3. What was the weakest part, and how would you fix it?",
            "4. Was the pace too fast / just right / too slow?",
            "5. Did theory and hands-on feel balanced?",
        ],
        size=17,
    )
    right = s2.shapes.add_textbox(Inches(6.9), Inches(1.35), Inches(5.9), Inches(5.8))
    add_bullets(
        right.text_frame,
        [
            "6. How clear and helpful were the facilitators?",
            "7. How useful were the materials (guides, .sb3, slides)?",
            "8. How confident do you feel to train/teach this after the workshop? (1–5)",
            "9. Would you recommend this workshop to a colleague? Why / why not?",
            "10. One change that would most improve quality next time?",
        ],
        size=17,
    )

    # Slide 3 — tips + exit ticket
    s3 = blank_slide(prs)
    add_title_bar(s3, "Facilitator Tips & Exit Ticket")
    box = s3.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tips = [
        ("Facilitator tips", True, 22, TEAL),
        ("Ask for examples, not slogans (“Day 2 Catch Game demo was clear”).", False, 18, DARK),
        ("Separate content vs delivery vs logistics.", False, 18, DARK),
        ("Don’t defend in the room — note it, thank them, move on.", False, 18, DARK),
        ("", False, 10, DARK),
        ("Exit ticket (write one sentence)", True, 22, TEAL),
        ("This workshop would be excellent if…", False, 20, DARK),
        ("", False, 10, DARK),
        ("Written form (optional detail): supporting-materials/workshop-evaluation-form.md", False, 15, MUTED),
        ("DSTI–CHPC Coding & Robotics Workshop", False, 14, MUTED),
    ]
    first = True
    for text, bold, size, color in tips:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        if not text:
            continue
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=bold, color=color)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
